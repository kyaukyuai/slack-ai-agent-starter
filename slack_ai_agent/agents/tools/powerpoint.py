"""PowerPoint generation tool."""

import os
from datetime import datetime
from typing import Any
from typing import Dict
from typing import List

import pytz  # type: ignore
from langchain.tools import StructuredTool


# Check if python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.util import Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print(
        "python-pptx package not installed. PowerPoint generation will not be available."
    )
    print("Install with: pip install python-pptx")


# Template slide layouts
TITLE_SLIDE_LAYOUT = 0
CONTENT_SLIDE_LAYOUT = 1
IMAGE_SLIDE_LAYOUT = 2
TABLE_SLIDE_LAYOUT = 3
TWO_COLUMN_SLIDE_LAYOUT = 4
THREE_IMAGES_SLIDE_LAYOUT = 5
THREE_HORIZONTAL_FLOW_SLIDE_LAYOUT = 6
THREE_VERTICAL_FLOW_SLIDE_LAYOUT = 7


def get_current_jst_time() -> str:
    """Get the current time in JST format.

    Returns:
        str: Current time in JST format (YYYYMMDD_HHMMSS)
    """
    jst = pytz.timezone("Asia/Tokyo")
    current_time = datetime.now(jst)
    return current_time.strftime("%Y%m%d_%H%M%S")


# Define a type alias for Presentation to handle the case when python-pptx is not installed

if PPTX_AVAILABLE:
    PresentationType = Presentation
else:
    PresentationType = Any  # type: ignore[assignment]


class SlideBase:
    """Base class for slide creation."""

    def __init__(self, presentation: Any, page: Dict):
        """Initialize the slide base.

        Args:
            presentation: The PowerPoint presentation object
            page: The page data containing content for the slide
        """
        self.presentation = presentation
        self.page = page
        self.slide = None

    def set_title(self, title: str) -> None:
        """Set the title of the slide.

        Args:
            title: The title text
        """
        if self.slide and self.slide.shapes.title:
            self.slide.shapes.title.text = title

    def add_sections_to_slide(self, text_frame, sections: List[Dict]) -> None:
        """Add sections to a slide.

        Args:
            text_frame: The text frame to add sections to
            sections: List of section data
        """
        for i, section in enumerate(sections):
            # Add section title
            p = text_frame.add_paragraph()
            p.text = section["title"]
            p.font.bold = True
            p.font.size = Pt(18)

            # Add section content (bullet points)
            for item in section["content"]:
                p = text_frame.add_paragraph()
                p.text = item
                p.level = 1
                p.font.size = Pt(14)

            # Add spacing between sections (except for the last one)
            if i < len(sections) - 1:
                text_frame.add_paragraph()


class TextSlide(SlideBase):
    """Text slide class."""

    def create_slide(self) -> None:
        """Create a text slide."""
        slide_layout = self.presentation.slide_layouts[CONTENT_SLIDE_LAYOUT]
        self.slide = self.presentation.slides.add_slide(slide_layout)
        self.set_title(self.page["header"])

        content = self.slide.placeholders[1]  # type: ignore
        text_frame = content.text_frame
        text_frame.clear()

        self.add_sections_to_slide(text_frame, self.page["sections"])


class ImageSlide(SlideBase):
    """Image slide class."""

    def create_slide(self) -> None:
        """Create an image slide."""
        slide_layout = self.presentation.slide_layouts[IMAGE_SLIDE_LAYOUT]
        self.slide = self.presentation.slides.add_slide(slide_layout)
        self.set_title(self.page["header"])

        content = self.slide.placeholders[1]  # type: ignore
        text_frame = content.text_frame
        text_frame.clear()

        self.add_sections_to_slide(text_frame, self.page["sections"])

        # Note: In a real implementation, we would add the image here
        # For now, we'll just add a placeholder note
        notes_slide = self.slide.notes_slide  # type: ignore
        notes_slide.notes_text_frame.text = (
            f"Image to be added: {self.page.get('image_path', 'No image specified')}"
        )


class TableSlide(SlideBase):
    """Table slide class."""

    def create_slide(self) -> None:
        """Create a table slide."""
        slide_layout = self.presentation.slide_layouts[TABLE_SLIDE_LAYOUT]
        self.slide = self.presentation.slides.add_slide(slide_layout)
        self.set_title(self.page["header"])

        # Add table
        table_data = self.page.get("table_data", {})
        rows = len(table_data.get("rows", [])) + 1  # +1 for header
        cols = len(table_data.get("headers", []))

        if rows > 1 and cols > 0:
            left = Inches(1.0)
            top = Inches(2.0)
            width = Inches(8.0)
            height = Inches(2.0)

            table = self.slide.shapes.add_table(  # type: ignore
                rows, cols, left, top, width, height
            ).table

            # Add headers
            for i, header in enumerate(table_data.get("headers", [])):
                cell = table.cell(0, i)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True

            # Add rows
            for i, row_data in enumerate(table_data.get("rows", [])):
                for j, cell_text in enumerate(row_data):
                    if j < cols:  # Ensure we don't exceed the number of columns
                        cell = table.cell(i + 1, j)  # +1 to skip header row
                        cell.text = cell_text

        # Add key messages
        if "key_messages" in self.page:
            content = self.slide.placeholders[1]  # type: ignore
            text_frame = content.text_frame
            text_frame.clear()

            p = text_frame.add_paragraph()
            p.text = "Key Messages:"
            p.font.bold = True

            for message in self.page["key_messages"]:
                p = text_frame.add_paragraph()
                p.text = message
                p.level = 1


class TwoColumnSlide(SlideBase):
    """Two column slide class."""

    def create_slide(self) -> None:
        """Create a two column slide."""
        slide_layout = self.presentation.slide_layouts[TWO_COLUMN_SLIDE_LAYOUT]
        self.slide = self.presentation.slides.add_slide(slide_layout)
        self.set_title(self.page["header"])

        # Left column
        left_content = self.slide.placeholders[1]  # type: ignore
        left_frame = left_content.text_frame
        left_frame.clear()

        p = left_frame.add_paragraph()
        p.text = self.page.get("left_title", "Left Column")
        p.font.bold = True

        for item in self.page.get("left_content", []):
            p = left_frame.add_paragraph()
            p.text = item
            p.level = 1

        # Right column
        right_content = self.slide.placeholders[2]  # type: ignore
        right_frame = right_content.text_frame
        right_frame.clear()

        p = right_frame.add_paragraph()
        p.text = self.page.get("right_title", "Right Column")
        p.font.bold = True

        for item in self.page.get("right_content", []):
            p = right_frame.add_paragraph()
            p.text = item
            p.level = 1


class ThreeImagesSlide(SlideBase):
    """Three images slide class."""

    def create_slide(self) -> None:
        """Create a slide with three images."""
        slide_layout = self.presentation.slide_layouts[THREE_IMAGES_SLIDE_LAYOUT]
        self.slide = self.presentation.slides.add_slide(slide_layout)
        self.set_title(self.page["header"])

        # For each image placeholder, add key message
        for i in range(3):
            if i < len(self.page.get("image_sections", [])):
                image_section = self.page["image_sections"][i]

                # Add key message
                content = self.slide.placeholders[  # type: ignore
                    i + 1
                ]  # +1 because placeholder 0 is the title
                text_frame = content.text_frame
                text_frame.clear()

                p = text_frame.add_paragraph()
                p.text = image_section.get("title", f"Image {i + 1}")
                p.font.bold = True

                for item in image_section.get("key_message", []):
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.level = 1

                # Note: In a real implementation, we would add the image here
                notes_slide = self.slide.notes_slide  # type: ignore
                notes_text = notes_slide.notes_text_frame.text
                notes_slide.notes_text_frame.text = f"{notes_text}\nImage {i + 1} to be added: {image_section.get('image_path', 'No image specified')}"


class ThreeHorizontalFlowSlide(SlideBase):
    """Three horizontal flow slide class."""

    def create_slide(self) -> None:
        """Create a slide with three horizontal flow steps."""
        slide_layout = self.presentation.slide_layouts[
            THREE_HORIZONTAL_FLOW_SLIDE_LAYOUT
        ]
        self.slide = self.presentation.slides.add_slide(slide_layout)
        self.set_title(self.page["header"])

        # For each step placeholder, add content
        for i in range(3):
            if i < len(self.page.get("steps", [])):
                step = self.page["steps"][i]

                content = self.slide.placeholders[  # type: ignore
                    i + 1
                ]  # +1 because placeholder 0 is the title
                text_frame = content.text_frame
                text_frame.clear()

                p = text_frame.add_paragraph()
                p.text = step.get("title", f"Step {i + 1}")
                p.font.bold = True

                for item in step.get("content", []):
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.level = 1


class ThreeVerticalFlowSlide(SlideBase):
    """Three vertical flow slide class."""

    def create_slide(self) -> None:
        """Create a slide with three vertical flow steps."""
        slide_layout = self.presentation.slide_layouts[THREE_VERTICAL_FLOW_SLIDE_LAYOUT]
        self.slide = self.presentation.slides.add_slide(slide_layout)
        self.set_title(self.page["header"])

        # For each step placeholder, add content
        for i in range(3):
            if i < len(self.page.get("steps", [])):
                step = self.page["steps"][i]

                content = self.slide.placeholders[  # type: ignore
                    i + 1
                ]  # +1 because placeholder 0 is the title
                text_frame = content.text_frame
                text_frame.clear()

                p = text_frame.add_paragraph()
                p.text = step.get("title", f"Step {i + 1}")
                p.font.bold = True

                for item in step.get("content", []):
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.level = 1


def create_powerpoint(title: str, slides: List[Dict]) -> str:
    """Create a PowerPoint presentation.

    Args:
        title: The title of the presentation
        slides: List of slide data

    Returns:
        str: Path to the created PowerPoint file
    """
    if not PPTX_AVAILABLE:
        return "Error: python-pptx package is not installed. Please install it with: pip install python-pptx"

    # Create a presentation with the default template
    prs = Presentation()

    # Add title slide
    title_slide_layout = prs.slide_layouts[TITLE_SLIDE_LAYOUT]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = f"作成日: {datetime.now().strftime('%Y/%m/%d')}"

    # Add content slides
    for slide_data in slides:
        template = slide_data.get("template", "text")

        if template == "text":
            TextSlide(prs, slide_data).create_slide()
        elif template == "image":
            ImageSlide(prs, slide_data).create_slide()
        elif template == "table":
            TableSlide(prs, slide_data).create_slide()
        elif template == "two_column":
            TwoColumnSlide(prs, slide_data).create_slide()
        elif template == "three_images":
            ThreeImagesSlide(prs, slide_data).create_slide()
        elif template == "three_horizontal_flow":
            ThreeHorizontalFlowSlide(prs, slide_data).create_slide()
        elif template == "three_vertical_flow":
            ThreeVerticalFlowSlide(prs, slide_data).create_slide()

    # Save the presentation
    output_dir = "output"
    os.makedirs(output_dir, exist_ok=True)

    timestamp = get_current_jst_time()
    filename = f"{title.replace(' ', '_')}_{timestamp}.pptx"
    file_path = os.path.join(output_dir, filename)

    prs.save(file_path)
    return file_path


def generate_powerpoint(topic: str, content: Dict) -> str:
    """Generate a PowerPoint presentation.

    Args:
        topic: The topic of the presentation
        content: Dictionary containing the presentation content

    Returns:
        str: Path to the created PowerPoint file or error message if python-pptx is not available
    """
    if not PPTX_AVAILABLE:
        return "Error: python-pptx package is not installed. Please install it with: pip install python-pptx"

    title = content.get("title", topic)
    slides = content.get("pages", [])

    return create_powerpoint(title, slides)


# Create the PowerPoint tool
create_powerpoint_tool = StructuredTool.from_function(
    func=generate_powerpoint,
    name="generate_powerpoint",
    description="Generate a PowerPoint presentation based on the provided topic and content. Returns the path to the created file.",
)
