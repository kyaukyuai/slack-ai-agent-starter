"""PowerPoint generation tool."""

import os
from datetime import datetime
from typing import Any
from typing import Dict
from typing import List

import pytz  # type: ignore
from langchain.tools import StructuredTool


# Check if python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.util import Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print(
        "python-pptx package not installed. PowerPoint generation will not be available."
    )
    print("Install with: pip install python-pptx")


# Template slide layouts
TITLE_SLIDE_LAYOUT = 0
CONTENT_SLIDE_LAYOUT = 1
IMAGE_SLIDE_LAYOUT = 2
TABLE_SLIDE_LAYOUT = 3
TWO_COLUMN_SLIDE_LAYOUT = 4
THREE_IMAGES_SLIDE_LAYOUT = 5
THREE_HORIZONTAL_FLOW_SLIDE_LAYOUT = 6
THREE_VERTICAL_FLOW_SLIDE_LAYOUT = 7


def get_current_jst_time() -> str:
    """Get the current time in JST format.

    Returns:
        str: Current time in JST format (YYYYMMDD_HHMMSS)
    """
    jst = pytz.timezone("Asia/Tokyo")
    current_time = datetime.now(jst)
    return current_time.strftime("%Y%m%d_%H%M%S")


class SlideBase:
    """Base class for slide creation."""

    def __init__(self, presentation: Any, page: Dict):
        """Initialize the slide base.

        Args:
            presentation: The PowerPoint presentation object
            page: The page data containing content for the slide
        """
        self.presentation = presentation
        self.page = page
        self.slide = None

    def set_title(self, title: str) -> None:
        """Set the title of the slide.

        Args:
            title: The title text
        """
        if self.slide and self.slide.shapes.title:
            self.slide.shapes.title.text = title

    def add_sections_to_slide(self, text_frame, sections: List[Dict]) -> None:
        """Add sections to a slide.

        Args:
            text_frame: The text frame to add sections to
            sections: List of section data
        """
        for i, section in enumerate(sections):
            # Add section title
            p = text_frame.add_paragraph()
            p.text = section["title"]
            p.font.bold = True
            p.font.size = Pt(18)

            # Add section content (bullet points)
            for item in section["content"]:
                p = text_frame.add_paragraph()
                p.text = item
                p.level = 1
                p.font.size = Pt(14)

            # Add spacing between sections (except for the last one)
            if i < len(sections) - 1:
                text_frame.add_paragraph()


class TextSlide(SlideBase):
    """Text slide class."""

    def create_slide(self) -> None:
        """Create a text slide."""
        try:
            slide_layout = self.presentation.slide_layouts[CONTENT_SLIDE_LAYOUT]
            self.slide = self.presentation.slides.add_slide(slide_layout)
            self.set_title(self.page["header"])

            # Try to get the content placeholder
            try:
                content = self.slide.placeholders[1]  # type: ignore
                text_frame = content.text_frame
                text_frame.clear()
            except (KeyError, IndexError):
                # If placeholder doesn't exist, create a text box
                left = Inches(1.0)
                top = Inches(2.0)
                width = Inches(8.0)
                height = Inches(4.0)
                shape = self.slide.shapes.add_textbox(left, top, width, height)  # type: ignore
                text_frame = shape.text_frame

            self.add_sections_to_slide(text_frame, self.page["sections"])
        except Exception as e:
            print(f"Error creating text slide: {str(e)}")
            # Create a simple slide with error message
            slide_layout = self.presentation.slide_layouts[0]  # Title slide as fallback
            self.slide = self.presentation.slides.add_slide(slide_layout)
            self.set_title(f"Error: {self.page['header']}")

            # Add error details
            left = Inches(1.0)
            top = Inches(2.0)
            width = Inches(8.0)
            height = Inches(4.0)
            shape = self.slide.shapes.add_textbox(left, top, width, height)  # type: ignore
            shape.text_frame.text = f"Error creating slide: {str(e)}"


class ImageSlide(SlideBase):
    """Image slide class."""

    def create_slide(self) -> None:
        """Create an image slide."""
        try:
            slide_layout = self.presentation.slide_layouts[IMAGE_SLIDE_LAYOUT]
            self.slide = self.presentation.slides.add_slide(slide_layout)
            self.set_title(self.page["header"])

            # Try to get the content placeholder
            try:
                content = self.slide.placeholders[1]  # type: ignore
                text_frame = content.text_frame
                text_frame.clear()
            except (KeyError, IndexError):
                # If placeholder doesn't exist, create a text box
                left = Inches(1.0)
                top = Inches(2.0)
                width = Inches(4.0)
                height = Inches(4.0)
                shape = self.slide.shapes.add_textbox(left, top, width, height)  # type: ignore
                text_frame = shape.text_frame

            self.add_sections_to_slide(text_frame, self.page["sections"])

            # Note: In a real implementation, we would add the image here
            # For now, we'll just add a placeholder note
            try:
                notes_slide = self.slide.notes_slide  # type: ignore
                notes_slide.notes_text_frame.text = f"Image to be added: {self.page.get('image_path', 'No image specified')}"
            except Exception as note_error:
                print(f"Error adding notes: {str(note_error)}")
        except Exception as e:
            print(f"Error creating image slide: {str(e)}")
            # Create a simple slide with error message
            slide_layout = self.presentation.slide_layouts[0]  # Title slide as fallback
            self.slide = self.presentation.slides.add_slide(slide_layout)
            self.set_title(f"Error: {self.page['header']}")

            # Add error details
            left = Inches(1.0)
            top = Inches(2.0)
            width = Inches(8.0)
            height = Inches(4.0)
            shape = self.slide.shapes.add_textbox(left, top, width, height)  # type: ignore
            shape.text_frame.text = f"Error creating slide: {str(e)}"


class TableSlide(SlideBase):
    """Table slide class."""

    def create_slide(self) -> None:
        """Create a table slide."""
        try:
            slide_layout = self.presentation.slide_layouts[TABLE_SLIDE_LAYOUT]
            self.slide = self.presentation.slides.add_slide(slide_layout)
            self.set_title(self.page["header"])

            # Add table
            table_data = self.page.get("table_data", {})
            rows = len(table_data.get("rows", [])) + 1  # +1 for header
            cols = len(table_data.get("headers", []))

            if rows > 1 and cols > 0:
                left = Inches(1.0)
                top = Inches(2.0)
                width = Inches(8.0)
                height = Inches(2.0)

                try:
                    table = self.slide.shapes.add_table(  # type: ignore
                        rows, cols, left, top, width, height
                    ).table  # type: ignore

                    # Add headers
                    for i, header in enumerate(table_data.get("headers", [])):
                        cell = table.cell(0, i)
                        cell.text = header
                        cell.text_frame.paragraphs[0].font.bold = True

                    # Add rows
                    for i, row_data in enumerate(table_data.get("rows", [])):
                        for j, cell_text in enumerate(row_data):
                            if j < cols:  # Ensure we don't exceed the number of columns
                                cell = table.cell(i + 1, j)  # +1 to skip header row
                                cell.text = cell_text
                except Exception as table_error:
                    print(f"Error creating table: {str(table_error)}")
                    # Add error message as text
                    left = Inches(1.0)
                    top = Inches(2.0)
                    width = Inches(8.0)
                    height = Inches(2.0)
                    shape = self.slide.shapes.add_textbox(left, top, width, height)  # type: ignore
                    shape.text_frame.text = f"Error creating table: {str(table_error)}"

            # Add key messages
            if "key_messages" in self.page:
                try:
                    content = get_placeholder_safely(self.slide, 1, "Key Messages:")
                    if hasattr(content, "text_frame"):
                        text_frame = content.text_frame
                        text_frame.clear()
                    else:
                        text_frame = (
                            content  # Already a text frame from get_placeholder_safely
                        )

                    p = text_frame.add_paragraph()
                    p.text = "Key Messages:"
                    p.font.bold = True

                    for message in self.page["key_messages"]:
                        p = text_frame.add_paragraph()
                        p.text = message
                        p.level = 1
                except Exception as key_error:
                    print(f"Error adding key messages: {str(key_error)}")
        except Exception as e:
            print(f"Error creating table slide: {str(e)}")
            # Create a simple slide with error message
            slide_layout = self.presentation.slide_layouts[0]  # Title slide as fallback
            self.slide = self.presentation.slides.add_slide(slide_layout)
            self.set_title(f"Error: {self.page['header']}")

            # Add error details
            left = Inches(1.0)
            top = Inches(2.0)
            width = Inches(8.0)
            height = Inches(4.0)
            shape = self.slide.shapes.add_textbox(left, top, width, height)  # type: ignore
            shape.text_frame.text = f"Error creating slide: {str(e)}"


class TwoColumnSlide(SlideBase):
    """Two column slide class."""

    def create_slide(self) -> None:
        """Create a two column slide."""
        try:
            slide_layout = self.presentation.slide_layouts[TWO_COLUMN_SLIDE_LAYOUT]
            self.slide = self.presentation.slides.add_slide(slide_layout)
            self.set_title(self.page["header"])

            # Left column
            try:
                left_content = get_placeholder_safely(self.slide, 1, "Left Column")
                if hasattr(left_content, "text_frame"):
                    left_frame = left_content.text_frame
                    left_frame.clear()
                else:
                    left_frame = (
                        left_content  # Already a text frame from get_placeholder_safely
                    )

                p = left_frame.add_paragraph()
                p.text = self.page.get("left_title", "Left Column")
                p.font.bold = True

                for item in self.page.get("left_content", []):
                    p = left_frame.add_paragraph()
                    p.text = item
                    p.level = 1
            except Exception as left_error:
                print(f"Error creating left column: {str(left_error)}")
                # Create a fallback text box
                left = Inches(0.5)
                top = Inches(2.0)
                width = Inches(4.0)
                height = Inches(4.0)
                shape = self.slide.shapes.add_textbox(left, top, width, height)  # type: ignore
                shape.text_frame.text = f"Left Column Error: {str(left_error)}"

            # Right column
            try:
                right_content = get_placeholder_safely(self.slide, 2, "Right Column")
                if hasattr(right_content, "text_frame"):
                    right_frame = right_content.text_frame
                    right_frame.clear()
                else:
                    right_frame = right_content  # Already a text frame from get_placeholder_safely

                p = right_frame.add_paragraph()
                p.text = self.page.get("right_title", "Right Column")
                p.font.bold = True

                for item in self.page.get("right_content", []):
                    p = right_frame.add_paragraph()
                    p.text = item
                    p.level = 1
            except Exception as right_error:
                print(f"Error creating right column: {str(right_error)}")
                # Create a fallback text box
                left = Inches(5.0)
                top = Inches(2.0)
                width = Inches(4.0)
                height = Inches(4.0)
                shape = self.slide.shapes.add_textbox(left, top, width, height)  # type: ignore
                shape.text_frame.text = f"Right Column Error: {str(right_error)}"
        except Exception as e:
            print(f"Error creating two column slide: {str(e)}")
            # Create a simple slide with error message
            slide_layout = self.presentation.slide_layouts[0]  # Title slide as fallback
            self.slide = self.presentation.slides.add_slide(slide_layout)
            self.set_title(f"Error: {self.page['header']}")

            # Add error details
            left = Inches(1.0)
            top = Inches(2.0)
            width = Inches(8.0)
            height = Inches(4.0)
            shape = self.slide.shapes.add_textbox(left, top, width, height)  # type: ignore
            shape.text_frame.text = f"Error creating slide: {str(e)}"


class ThreeImagesSlide(SlideBase):
    """Three images slide class."""

    def create_slide(self) -> None:
        """Create a slide with three images."""
        try:
            slide_layout = self.presentation.slide_layouts[THREE_IMAGES_SLIDE_LAYOUT]
            self.slide = self.presentation.slides.add_slide(slide_layout)
            self.set_title(self.page["header"])

            # For each image placeholder, add key message
            for i in range(3):
                if i < len(self.page.get("image_sections", [])):
                    image_section = self.page["image_sections"][i]

                    try:
                        # Add key message
                        content = get_placeholder_safely(
                            self.slide, i + 1, f"Image {i + 1}"
                        )  # +1 because placeholder 0 is the title

                        if hasattr(content, "text_frame"):
                            text_frame = content.text_frame
                            text_frame.clear()
                        else:
                            text_frame = content  # Already a text frame from get_placeholder_safely

                        p = text_frame.add_paragraph()
                        p.text = image_section.get("title", f"Image {i + 1}")
                        p.font.bold = True

                        for item in image_section.get("key_message", []):
                            p = text_frame.add_paragraph()
                            p.text = item
                            p.level = 1

                        # Note: In a real implementation, we would add the image here
                        try:
                            notes_slide = self.slide.notes_slide  # type: ignore
                            notes_text = (
                                notes_slide.notes_text_frame.text
                                if hasattr(notes_slide.notes_text_frame, "text")
                                else ""
                            )
                            notes_slide.notes_text_frame.text = f"{notes_text}\nImage {i + 1} to be added: {image_section.get('image_path', 'No image specified')}"
                        except Exception as note_error:
                            print(
                                f"Error adding notes for image {i + 1}: {str(note_error)}"
                            )
                    except Exception as image_error:
                        print(f"Error adding image section {i + 1}: {str(image_error)}")
                        # Create a fallback text box for this image
                        left = Inches(1.0 + i * 2.5)
                        top = Inches(2.0)
                        width = Inches(2.0)
                        height = Inches(3.0)
                        shape = self.slide.shapes.add_textbox(left, top, width, height)  # type: ignore
                        shape.text_frame.text = (
                            f"Image {i + 1} Error: {str(image_error)}"
                        )
        except Exception as e:
            print(f"Error creating three images slide: {str(e)}")
            # Create a simple slide with error message
            slide_layout = self.presentation.slide_layouts[0]  # Title slide as fallback
            self.slide = self.presentation.slides.add_slide(slide_layout)
            self.set_title(f"Error: {self.page['header']}")

            # Add error details
            left = Inches(1.0)
            top = Inches(2.0)
            width = Inches(8.0)
            height = Inches(4.0)
            shape = self.slide.shapes.add_textbox(left, top, width, height)  # type: ignore
            shape.text_frame.text = f"Error creating slide: {str(e)}"


class ThreeHorizontalFlowSlide(SlideBase):
    """Three horizontal flow slide class."""

    def create_slide(self) -> None:
        """Create a slide with three horizontal flow steps."""
        try:
            slide_layout = self.presentation.slide_layouts[
                THREE_HORIZONTAL_FLOW_SLIDE_LAYOUT
            ]
            self.slide = self.presentation.slides.add_slide(slide_layout)
            self.set_title(self.page["header"])

            # For each step placeholder, add content
            for i in range(3):
                if i < len(self.page.get("steps", [])):
                    step = self.page["steps"][i]

                    try:
                        content = get_placeholder_safely(
                            self.slide, i + 1, f"Step {i + 1}"
                        )  # +1 because placeholder 0 is the title

                        if hasattr(content, "text_frame"):
                            text_frame = content.text_frame
                            text_frame.clear()
                        else:
                            text_frame = content  # Already a text frame from get_placeholder_safely

                        p = text_frame.add_paragraph()
                        p.text = step.get("title", f"Step {i + 1}")
                        p.font.bold = True

                        for item in step.get("content", []):
                            p = text_frame.add_paragraph()
                            p.text = item
                            p.level = 1
                    except Exception as step_error:
                        print(f"Error adding step {i + 1}: {str(step_error)}")
                        # Create a fallback text box for this step
                        left = Inches(0.5 + i * 3.0)
                        top = Inches(2.0)
                        width = Inches(2.5)
                        height = Inches(3.0)
                        shape = self.slide.shapes.add_textbox(left, top, width, height)  # type: ignore
                        shape.text_frame.text = f"Step {i + 1} Error: {str(step_error)}"
        except Exception as e:
            print(f"Error creating horizontal flow slide: {str(e)}")
            # Create a simple slide with error message
            slide_layout = self.presentation.slide_layouts[0]  # Title slide as fallback
            self.slide = self.presentation.slides.add_slide(slide_layout)
            self.set_title(f"Error: {self.page['header']}")

            # Add error details
            left = Inches(1.0)
            top = Inches(2.0)
            width = Inches(8.0)
            height = Inches(4.0)
            shape = self.slide.shapes.add_textbox(left, top, width, height)  # type: ignore
            shape.text_frame.text = f"Error creating slide: {str(e)}"


class ThreeVerticalFlowSlide(SlideBase):
    """Three vertical flow slide class."""

    def create_slide(self) -> None:
        """Create a slide with three vertical flow steps."""
        try:
            slide_layout = self.presentation.slide_layouts[
                THREE_VERTICAL_FLOW_SLIDE_LAYOUT
            ]
            self.slide = self.presentation.slides.add_slide(slide_layout)
            self.set_title(self.page["header"])

            # For each step placeholder, add content
            for i in range(3):
                if i < len(self.page.get("steps", [])):
                    step = self.page["steps"][i]

                    try:
                        content = get_placeholder_safely(
                            self.slide, i + 1, f"Step {i + 1}"
                        )  # +1 because placeholder 0 is the title

                        if hasattr(content, "text_frame"):
                            text_frame = content.text_frame
                            text_frame.clear()
                        else:
                            text_frame = content  # Already a text frame from get_placeholder_safely

                        p = text_frame.add_paragraph()
                        p.text = step.get("title", f"Step {i + 1}")
                        p.font.bold = True

                        for item in step.get("content", []):
                            p = text_frame.add_paragraph()
                            p.text = item
                            p.level = 1
                    except Exception as step_error:
                        print(f"Error adding step {i + 1}: {str(step_error)}")
                        # Create a fallback text box for this step
                        left = Inches(1.0)
                        top = Inches(2.0 + i * 1.5)
                        width = Inches(8.0)
                        height = Inches(1.0)
                        shape = self.slide.shapes.add_textbox(left, top, width, height)  # type: ignore
                        shape.text_frame.text = f"Step {i + 1} Error: {str(step_error)}"
        except Exception as e:
            print(f"Error creating vertical flow slide: {str(e)}")
            # Create a simple slide with error message
            slide_layout = self.presentation.slide_layouts[0]  # Title slide as fallback
            self.slide = self.presentation.slides.add_slide(slide_layout)
            self.set_title(f"Error: {self.page['header']}")

            # Add error details
            left = Inches(1.0)
            top = Inches(2.0)
            width = Inches(8.0)
            height = Inches(4.0)
            shape = self.slide.shapes.add_textbox(left, top, width, height)  # type: ignore
            shape.text_frame.text = f"Error creating slide: {str(e)}"


def get_placeholder_safely(slide, idx, default_text=""):
    """Get a placeholder safely, returning None if it doesn't exist.

    Args:
        slide: The slide to get the placeholder from
        idx: The index of the placeholder
        default_text: Default text to use if placeholder doesn't exist

    Returns:
        The placeholder if it exists, or None
    """
    try:
        return slide.placeholders[idx]
    except (KeyError, IndexError):
        # If placeholder doesn't exist, create a text box instead
        left = Inches(1.0)
        top = Inches(2.0)
        width = Inches(8.0)
        height = Inches(4.0)

        # Add a text box with the default text
        txBox = slide.shapes.add_textbox(left, top, width, height)  # type: ignore
        tf = txBox.text_frame
        tf.text = default_text

        return tf


def create_powerpoint(title: str, slides: List[Dict]) -> str:
    """Create a PowerPoint presentation.

    Args:
        title: The title of the presentation
        slides: List of slide data

    Returns:
        str: Path to the created PowerPoint file
    """
    if not PPTX_AVAILABLE:
        return "Error: python-pptx package is not installed. Please install it with: pip install python-pptx"

    # Create a presentation with the default template
    prs = Presentation()

    # Add title slide
    title_slide_layout = prs.slide_layouts[TITLE_SLIDE_LAYOUT]
    slide = prs.slides.add_slide(title_slide_layout)

    # Safely get title and subtitle
    if slide.shapes.title:
        slide.shapes.title.text = title

    try:
        subtitle_shape = slide.placeholders[1]
        if hasattr(subtitle_shape, "text_frame"):
            subtitle_shape.text_frame.text = (
                f"作成日: {datetime.now().strftime('%Y/%m/%d')}"
            )
        else:
            # Try direct text assignment if text_frame is not available
            try:
                subtitle_shape.text = f"作成日: {datetime.now().strftime('%Y/%m/%d')}"
            except AttributeError:
                # If neither works, create a text box
                left = Inches(1.0)
                top = Inches(2.0)
                width = Inches(8.0)
                height = Inches(1.0)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                txBox.text_frame.text = f"作成日: {datetime.now().strftime('%Y/%m/%d')}"
    except (KeyError, IndexError):
        # If subtitle placeholder doesn't exist, add it as a text box
        left = Inches(1.0)
        top = Inches(2.0)
        width = Inches(8.0)
        height = Inches(1.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.text_frame.text = f"作成日: {datetime.now().strftime('%Y/%m/%d')}"

    # Add content slides
    for slide_data in slides:
        try:
            template = slide_data.get("template", "text")

            if template == "text":
                TextSlide(prs, slide_data).create_slide()
            elif template == "image":
                ImageSlide(prs, slide_data).create_slide()
            elif template == "table":
                TableSlide(prs, slide_data).create_slide()
            elif template == "two_column":
                TwoColumnSlide(prs, slide_data).create_slide()
            elif template == "three_images":
                ThreeImagesSlide(prs, slide_data).create_slide()
            elif template == "three_horizontal_flow":
                ThreeHorizontalFlowSlide(prs, slide_data).create_slide()
            elif template == "three_vertical_flow":
                ThreeVerticalFlowSlide(prs, slide_data).create_slide()
        except Exception as e:
            # If there's an error creating a slide, add a simple text slide with the error
            print(f"Error creating slide: {str(e)}")
            simple_layout = prs.slide_layouts[0]  # Use title slide layout as fallback
            error_slide = prs.slides.add_slide(simple_layout)
            if error_slide.shapes.title:
                error_slide.shapes.title.text = (
                    f"Error: {slide_data.get('header', 'Slide Error')}"
                )

            # Add error details as a text box
            left = Inches(1.0)
            top = Inches(2.0)
            width = Inches(8.0)
            height = Inches(4.0)
            txBox = error_slide.shapes.add_textbox(left, top, width, height)
            txBox.text_frame.text = (
                f"Error creating slide: {str(e)}\n\nSlide data: {str(slide_data)}"
            )

    # Save the presentation
    output_dir = "output"
    os.makedirs(output_dir, exist_ok=True)

    timestamp = get_current_jst_time()
    filename = f"{title.replace(' ', '_')}_{timestamp}.pptx"
    file_path = os.path.join(output_dir, filename)

    prs.save(file_path)
    return file_path


def generate_powerpoint(topic: str, content: Dict) -> str:
    """Generate a PowerPoint presentation.

    Args:
        topic: The topic of the presentation
        content: Dictionary containing the presentation content

    Returns:
        str: Path to the created PowerPoint file or error message if python-pptx is not available
    """
    if not PPTX_AVAILABLE:
        return "Error: python-pptx package is not installed. Please install it with: pip install python-pptx"

    title = content.get("title", topic)
    slides = content.get("pages", [])

    return create_powerpoint(title, slides)


# Create the PowerPoint tool
create_powerpoint_tool = StructuredTool.from_function(
    func=generate_powerpoint,
    name="generate_powerpoint",
    description="Generate a PowerPoint presentation based on the provided topic and content. Returns the path to the created file.",
)
